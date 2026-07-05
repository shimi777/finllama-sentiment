"""
Stage 2 evidence extraction: dump all slide text (including tables and speaker
notes) from BOTH implementation decks, extract quantitative claims into a CSV,
and diff the two decks for any NUMBER that differs between them.

Inputs (read directly -- decks are not Word-locked, no copy needed):
  presentation/implementation_deck.pptx
  presentation/implementation_deck_2026-05-04-07-48.pptx

Outputs:
  review/evidence/deck_main_dump.md    slide-by-slide text dump (main deck)
  review/evidence/deck_dated_dump.md   slide-by-slide text dump (dated deck)
  review/evidence/deck_claims.csv      one row per quantitative claim, both decks
  review/evidence/deck_diff.md         text diff + explicit number-level diff

Run:
  C:/python_projects/finllama-sentiment/.venv/Scripts/python.exe review/scripts/extract_pptx_text.py
"""
import csv
import difflib
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

REPO_ROOT = Path(__file__).resolve().parents[2]
MAIN_PPTX = REPO_ROOT / "presentation" / "implementation_deck.pptx"
DATED_PPTX = REPO_ROOT / "presentation" / "implementation_deck_2026-05-04-07-48.pptx"

DUMP_MAIN = REPO_ROOT / "review" / "evidence" / "deck_main_dump.md"
DUMP_DATED = REPO_ROOT / "review" / "evidence" / "deck_dated_dump.md"
CSV_PATH = REPO_ROOT / "review" / "evidence" / "deck_claims.csv"
DIFF_PATH = REPO_ROOT / "review" / "evidence" / "deck_diff.md"

# Reuse the same number-matching approach as the docx extractor.
RE_DECIMAL = re.compile(r"(?<![\w.])\d{1,4}\.\d+(?![\w])")
RE_PERCENT = re.compile(r"(?<![\w.])\d{1,3}(?:\.\d+)?\s?%")
RE_POINT_GAP = re.compile(r"(?<![\w.])\d{1,3}(?:\.\d+)?[\s-]?(?:point|pt)s?\b", re.IGNORECASE)
RE_MONEY = re.compile(r"[$₪€]\s?\d[\d,]*(?:\.\d+)?")
RE_TIME = re.compile(
    r"(?<![\w.])\d+(?:\.\d+)?\s?(?:ms|milliseconds?|s|sec|secs?|seconds?|min|mins?|minutes?|hrs?|hours?|days?)\b",
    re.IGNORECASE,
)
RE_PARAM_SIZE = re.compile(r"(?<![\w.])\d{1,4}(?:\.\d+)?\s?[MB](?![\w])")
RE_COUNT = re.compile(r"(?<![\w.,])\d{2,6}(?![\w.])")
RE_YEAR = re.compile(r"(?<!\d)(19|20)\d{2}(?!\d)")
# any numeric token at all, for the deck-vs-deck number diff
RE_ANY_NUMBER = re.compile(r"(?<![\w.])\d[\d,]*\.?\d*%?")

MODEL_KEYWORDS = [
    "FinLLaMA", "LLaMA-3.1-8B", "Llama-3.1-8B", "LLaMA", "Llama", "FinBERT", "VADER",
    "GPT-4", "GPT-3.5", "GPT-4o", "Instruct", "plutus", "Mistral", "Qwen", "GLiNER",
]
DATASET_KEYWORDS = ["FPB", "FiQA", "FiQA-SA", "Financial PhraseBank", "financial_phrasebank", "FiNER-ORD", "GLiNER", "FIN", "Alvarado"]
TEMPLATE_KEYWORDS = ["zero-shot", "few-shot", "zero shot", "few shot", "template", "CoT", "chain-of-thought", "shot"]
METRIC_KEYWORDS = {
    "f1_macro": ["f1", "f1-macro", "f1 macro", "macro-f1", "macro f1"],
    "accuracy": ["accuracy", "acc."],
    "coverage": ["coverage", "parse rate", "parse failure", "parse-fail"],
}


def infer_context_keys(sentence: str) -> str:
    keys = []
    low = sentence.lower()
    for kw in MODEL_KEYWORDS:
        if kw.lower() in low and kw not in keys:
            keys.append(kw)
    for kw in DATASET_KEYWORDS:
        if kw.lower() in low and kw not in keys:
            keys.append(kw)
    for kw in TEMPLATE_KEYWORDS:
        if kw.lower() in low and kw not in keys:
            keys.append(kw)
    return "; ".join(keys)


def infer_metric_type(sentence: str, raw_value: str) -> str:
    low = sentence.lower()
    for mtype, kws in METRIC_KEYWORDS.items():
        for kw in kws:
            if kw in low:
                return mtype
    if "%" in raw_value:
        return "percent"
    if re.search(r"[$₪€]", raw_value):
        return "money"
    if RE_POINT_GAP.search(raw_value):
        return "percent"
    if RE_TIME.search(raw_value):
        return "time"
    if RE_PARAM_SIZE.fullmatch(raw_value.strip()):
        return "count"
    if RE_DECIMAL.fullmatch(raw_value.strip()):
        try:
            v = float(raw_value)
            if 0.0 <= v <= 1.0:
                return "f1_macro"
        except ValueError:
            pass
        return "other"
    if raw_value.strip().isdigit():
        return "count"
    return "other"


def is_year_context(match_text: str) -> bool:
    return bool(RE_YEAR.fullmatch(match_text))


def is_citation_or_figure_or_section(sentence: str, match_start: int, match_text: str) -> bool:
    window_start = max(0, match_start - 25)
    prefix = sentence[window_start:match_start].lower()
    if prefix.rstrip().endswith("["):
        if "." not in match_text:
            return True
        return False
    exclude_cues = [
        "section ", "§", "figure ", "fig. ", "fig ", "table ", "chapter ",
        "appendix ", "eq. ", "equation ", "ref. ", "page ", "p. ", "slide ",
    ]
    for cue in exclude_cues:
        if prefix.rstrip().endswith(cue.rstrip()) or cue in prefix[-12:]:
            return True
    return False


def is_bare_slide_number(sentence: str, match_start: int, match_end: int, match_text: str) -> bool:
    """Exclude a standalone slide-number token: the ENTIRE text box content is just the number."""
    return sentence.strip() == match_text.strip()


def extract_numeric_claims(sentence: str):
    spans_found = []
    for regex in (RE_MONEY, RE_POINT_GAP, RE_TIME, RE_PERCENT, RE_PARAM_SIZE, RE_DECIMAL, RE_COUNT):
        for m in regex.finditer(sentence):
            spans_found.append((m.start(), m.end(), m.group()))

    spans_found.sort(key=lambda t: (t[0], -(t[1] - t[0])))
    kept = []
    for s, e, txt in spans_found:
        overlap = False
        for ks, ke, _ in kept:
            if s < ke and e > ks:
                overlap = True
                break
        if not overlap:
            kept.append((s, e, txt))

    results = []
    for s, e, txt in kept:
        if is_year_context(txt):
            continue
        if is_citation_or_figure_or_section(sentence, s, txt):
            continue
        metric_type = infer_metric_type(sentence, txt)
        results.append((txt.strip(), metric_type))
    return results


# ---------------------------------------------------------------------------
# PPTX walk
# ---------------------------------------------------------------------------

def iter_table_cells(table):
    """Yield (row_idx, col_idx, row_header, col_header, cell_text)."""
    grid = []
    for row in table.rows:
        grid.append([cell.text.strip() for cell in row.cells])
    if not grid:
        return
    header_row = grid[0]
    for r_idx, row in enumerate(grid):
        row_header = row[0] if row else ""
        for c_idx, cell_text in enumerate(row):
            col_header = header_row[c_idx] if c_idx < len(header_row) else f"col{c_idx}"
            yield r_idx, c_idx, row_header, col_header, cell_text


def dump_and_collect(pptx_path: Path, dump_path: Path, deck_label: str, claim_id_prefix: str, claim_counter_start=1):
    prs = Presentation(str(pptx_path))
    dump_lines = [f"# Deck dump: {pptx_path.name}\n"]
    claims = []
    counter = claim_counter_start

    for slide_idx, slide in enumerate(prs.slides, start=1):
        dump_lines.append(f"\n## Slide {slide_idx}\n")
        slide_title = None

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if not text:
                    continue
                if slide_title is None and shape == (slide.shapes.title if slide.shapes.title else None):
                    slide_title = text
                dump_lines.append(f"- {text}")

                sentences = re.split(r"(?<=[.!?])\s+(?=[A-Z(])|\n", text)
                if not sentences:
                    sentences = [text]
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                    if is_bare_slide_number(sentence, 0, len(sentence), sentence) and sentence.isdigit() and len(sentence) <= 2:
                        continue  # bare slide-number textbox
                    found = extract_numeric_claims(sentence)
                    for raw_value, metric_type in found:
                        claim_id = f"{claim_id_prefix}{counter:03d}"
                        counter += 1
                        claims.append({
                            "claim_id": claim_id,
                            "deck": deck_label,
                            "section": f"Slide {slide_idx}" + (f" — {slide_title}" if slide_title and slide_title != sentence else ""),
                            "sentence": sentence,
                            "metric_type": metric_type,
                            "value": raw_value,
                            "context_keys": infer_context_keys(sentence),
                        })

            if shape.has_table:
                table = shape.table
                dump_lines.append(f"\n**[Table — Slide {slide_idx}]**\n")
                grid = [[c.text.strip() for c in row.cells] for row in table.rows]
                if grid:
                    dump_lines.append("| " + " | ".join(grid[0]) + " |")
                    dump_lines.append("|" + "---|" * len(grid[0]))
                    for row in grid[1:]:
                        dump_lines.append("| " + " | ".join(row) + " |")

                for r_idx, c_idx, row_header, col_header, cell_text in iter_table_cells(table):
                    if r_idx == 0:
                        continue  # header row itself, skip claim extraction (unless it has numbers meaningfully, rare)
                    if not cell_text:
                        continue
                    cell_context = (
                        f"[Table, Slide {slide_idx}] row='{row_header}' col='{col_header}' cell='{cell_text}'"
                    )
                    found = extract_numeric_claims(cell_text)
                    combined_context = f"{row_header} {col_header} {cell_text}".lower()
                    for raw_value, _ in found:
                        refined_metric = infer_metric_type(combined_context, raw_value)
                        claim_id = f"{claim_id_prefix}{counter:03d}"
                        counter += 1
                        claims.append({
                            "claim_id": claim_id,
                            "deck": deck_label,
                            "section": f"Slide {slide_idx}",
                            "sentence": cell_context,
                            "metric_type": refined_metric,
                            "value": raw_value,
                            "context_keys": infer_context_keys(f"{row_header} {col_header}"),
                        })

        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""
            if notes_text:
                dump_lines.append(f"\n_Speaker notes (Slide {slide_idx}):_ {notes_text}\n")
                sentences = re.split(r"(?<=[.!?])\s+(?=[A-Z(])|\n", notes_text)
                for sentence in sentences:
                    sentence = sentence.strip()
                    if not sentence:
                        continue
                    found = extract_numeric_claims(sentence)
                    for raw_value, metric_type in found:
                        claim_id = f"{claim_id_prefix}{counter:03d}"
                        counter += 1
                        claims.append({
                            "claim_id": claim_id,
                            "deck": deck_label,
                            "section": f"Slide {slide_idx} (speaker notes)",
                            "sentence": sentence,
                            "metric_type": metric_type,
                            "value": raw_value,
                            "context_keys": infer_context_keys(sentence),
                        })

    dump_path.write_text("\n".join(dump_lines), encoding="utf-8")
    return claims, counter


def get_all_slide_text_lines(pptx_path: Path):
    """Flat list of (slide_idx, source, text) for diffing purposes."""
    prs = Presentation(str(pptx_path))
    lines = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append((slide_idx, "shape", text))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            lines.append((slide_idx, "table_cell", t))
        if slide.has_notes_slide:
            nt = slide.notes_slide.notes_text_frame.text.strip() if slide.notes_slide.notes_text_frame else ""
            if nt:
                lines.append((slide_idx, "notes", nt))
    return lines


def get_slide_tables(pptx_path: Path):
    """dict: slide_idx -> list of tables, each table a list-of-lists of cell text."""
    prs = Presentation(str(pptx_path))
    out = {}
    for slide_idx, slide in enumerate(prs.slides, start=1):
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                grid = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
                tables.append(grid)
        if tables:
            out[slide_idx] = tables
    return out


def build_table_cell_diff(main_tables, dated_tables):
    """Compare same-position table cells slide-by-slide (best-effort by slide index
    and table/row/col position) and report every cell whose NUMBER content differs,
    including 'filled in' cases where one side is a placeholder like [NUM] and the
    other has a concrete value.
    """
    diffs = []
    all_slides = sorted(set(main_tables) | set(dated_tables))
    for slide_idx in all_slides:
        m_tables = main_tables.get(slide_idx, [])
        d_tables = dated_tables.get(slide_idx, [])
        for t_idx in range(max(len(m_tables), len(d_tables))):
            m_grid = m_tables[t_idx] if t_idx < len(m_tables) else []
            d_grid = d_tables[t_idx] if t_idx < len(d_tables) else []
            max_rows = max(len(m_grid), len(d_grid))
            for r in range(max_rows):
                m_row = m_grid[r] if r < len(m_grid) else []
                d_row = d_grid[r] if r < len(d_grid) else []
                max_cols = max(len(m_row), len(d_row))
                for c in range(max_cols):
                    m_cell = m_row[c] if c < len(m_row) else "(missing)"
                    d_cell = d_row[c] if c < len(d_row) else "(missing)"
                    if m_cell == d_cell:
                        continue
                    m_nums = RE_ANY_NUMBER.findall(m_cell)
                    d_nums = RE_ANY_NUMBER.findall(d_cell)
                    has_placeholder = "[NUM]" in m_cell or "[NUM]" in d_cell
                    if m_nums != d_nums or has_placeholder:
                        row_header = m_row[0] if m_row else (d_row[0] if d_row else "")
                        col_header = (m_grid[0][c] if m_grid and c < len(m_grid[0]) else
                                      (d_grid[0][c] if d_grid and c < len(d_grid[0]) else f"col{c}"))
                        diffs.append({
                            "slide": slide_idx,
                            "table": t_idx + 1,
                            "row": r,
                            "col": c,
                            "row_header": row_header,
                            "col_header": col_header,
                            "main_cell": m_cell,
                            "dated_cell": d_cell,
                        })
    return diffs


def build_number_diff(main_lines, dated_lines):
    """Find slide-aligned lines whose text differs only/also in a NUMBER."""
    main_by_slide = {}
    for idx, src, text in main_lines:
        main_by_slide.setdefault(idx, []).append((src, text))
    dated_by_slide = {}
    for idx, src, text in dated_lines:
        dated_by_slide.setdefault(idx, []).append((src, text))

    diff_entries = []
    all_slides = sorted(set(main_by_slide) | set(dated_by_slide))
    for slide_idx in all_slides:
        main_texts = [t for _, t in main_by_slide.get(slide_idx, [])]
        dated_texts = [t for _, t in dated_by_slide.get(slide_idx, [])]

        # try to pair up similar lines between the two slides using difflib
        sm = difflib.SequenceMatcher(a=main_texts, b=dated_texts)
        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            if tag == "equal":
                continue
            main_chunk = main_texts[i1:i2]
            dated_chunk = dated_texts[j1:j2]
            # pair line-by-line best-effort
            for m_text in main_chunk:
                best_match = None
                best_ratio = 0.0
                for d_text in dated_chunk:
                    ratio = difflib.SequenceMatcher(a=m_text, b=d_text).ratio()
                    if ratio > best_ratio:
                        best_ratio = ratio
                        best_match = d_text
                if best_match is not None and best_ratio > 0.5:
                    main_nums = RE_ANY_NUMBER.findall(m_text)
                    dated_nums = RE_ANY_NUMBER.findall(best_match)
                    if main_nums != dated_nums:
                        diff_entries.append({
                            "slide": slide_idx,
                            "main_text": m_text,
                            "dated_text": best_match,
                            "main_numbers": main_nums,
                            "dated_numbers": dated_nums,
                        })
    return diff_entries


def main():
    if not MAIN_PPTX.exists() or not DATED_PPTX.exists():
        print("ERROR: one or both deck files not found.")
        sys.exit(1)

    main_claims, next_counter = dump_and_collect(MAIN_PPTX, DUMP_MAIN, "main", "D")
    dated_claims, _ = dump_and_collect(DATED_PPTX, DUMP_DATED, "dated", "D", claim_counter_start=next_counter)

    all_claims = main_claims + dated_claims
    with open(CSV_PATH, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(
            f, fieldnames=["claim_id", "deck", "section", "sentence", "metric_type", "value", "context_keys"]
        )
        writer.writeheader()
        for row in all_claims:
            writer.writerow(row)

    print(f"Wrote {DUMP_MAIN}")
    print(f"Wrote {DUMP_DATED}")
    print(f"Wrote {CSV_PATH} ({len(all_claims)} claims: {len(main_claims)} main, {len(dated_claims)} dated)")

    # ---- Deck diff ----
    main_lines = get_all_slide_text_lines(MAIN_PPTX)
    dated_lines = get_all_slide_text_lines(DATED_PPTX)

    main_all_text = "\n".join(f"[Slide {i}] {t}" for i, _, t in main_lines)
    dated_all_text = "\n".join(f"[Slide {i}] {t}" for i, _, t in dated_lines)

    text_diff = list(difflib.unified_diff(
        main_all_text.splitlines(), dated_all_text.splitlines(),
        fromfile="implementation_deck.pptx", tofile="implementation_deck_2026-05-04-07-48.pptx",
        lineterm="",
    ))

    number_diffs = build_number_diff(main_lines, dated_lines)

    main_tables = get_slide_tables(MAIN_PPTX)
    dated_tables = get_slide_tables(DATED_PPTX)
    table_cell_diffs = build_table_cell_diff(main_tables, dated_tables)

    diff_md = ["# Deck diff: main vs dated\n"]
    diff_md.append(f"- Main deck slides: {max((i for i,_,_ in main_lines), default=0)}")
    diff_md.append(f"- Dated deck slides: {max((i for i,_,_ in dated_lines), default=0)}\n")

    diff_md.append("## Table-cell number differences (position-aligned by slide/table/row/col)\n")
    if table_cell_diffs:
        for d in table_cell_diffs:
            diff_md.append(
                f"### Slide {d['slide']}, Table {d['table']}, row='{d['row_header']}', col='{d['col_header']}'"
            )
            diff_md.append(f"- MAIN:  `{d['main_cell']}`")
            diff_md.append(f"- DATED: `{d['dated_cell']}`")
            diff_md.append("")
    else:
        diff_md.append("(no table-cell differences found)\n")

    diff_md.append("## Numbers that differ between decks (slide-aligned best-effort line match, non-table text)\n")
    if number_diffs:
        for d in number_diffs:
            diff_md.append(f"### Slide {d['slide']}")
            diff_md.append(f"- MAIN:  `{d['main_text']}`  numbers={d['main_numbers']}")
            diff_md.append(f"- DATED: `{d['dated_text']}`  numbers={d['dated_numbers']}")
            diff_md.append("")
    else:
        diff_md.append("(none found by the slide-aligned line matcher)\n")

    diff_md.append("## Full unified text diff (raw)\n")
    diff_md.append("```diff")
    diff_md.extend(text_diff)
    diff_md.append("```")

    DIFF_PATH.write_text("\n".join(diff_md), encoding="utf-8")
    print(f"Wrote {DIFF_PATH} ({len(number_diffs)} number-level differences found)")


if __name__ == "__main__":
    main()
